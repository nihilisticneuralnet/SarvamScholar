import os
import streamlit as st
import google.generativeai as genai
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import requests
import base64
import json
import io
from IPython.display import Audio
import subprocess
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

load_dotenv()

os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

SARVAM_API_KEY = os.getenv("SARVAM_API_KEY")
SARVAM_API_URL = "https://api.sarvam.ai"

def create_ppt(module_content, module_name, is_hindi=False):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Module: {module_name}"
    subtitle.text = "Generated Content"
    
    content_slide_layout = prs.slide_layouts[1]
    
    paragraphs = module_content.split('\n\n')
    for para in paragraphs:
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = module_name
        content.text = para
    
    # Add a slide for the quiz if available
    if module_name in st.session_state.quizzes:
        quiz_slide = prs.slides.add_slide(content_slide_layout)
        quiz_title = quiz_slide.shapes.title
        quiz_content = quiz_slide.placeholders[1]
        
        quiz_title.text = f"Quiz: {module_name}"
        quiz_content.text = st.session_state.quizzes[module_name]
    
    output = BytesIO()
    prs.save(output)
    return output.getvalue()

def generate_module_content(title, module_name, course_info, previous_module=None):
    input_text = f"""
    Course Title: {title}
    Module Name: {module_name}
    Profession: {course_info['profession']}
    Target Audience: {course_info['target_audience']}
    Course Content: {course_info['course_content']}
    Course Description: {course_info['course_description']}
    Previous Module: {previous_module if previous_module else 'None'}
    """
    prompt = """
    Based on the following information, generate lecture content for the specified module:
    {input}

    Start the content with a brief recap of the previous module (if any) and an introduction to the current module.
    Use proper headings and subheadings to structure the content.
    The content should be educational, engaging, and relevant to the course and target audience.
    Ensure a smooth transition from the previous module's content to the current one.

    Format the content as follows:
    # Module Title

    ## Introduction
    [Brief recap of previous module if applicable]
    [Introduction to current module]

    ## [Main Topic 1]
    [Content]

    ### [Subtopic 1.1]
    [Content]

    ### [Subtopic 1.2]
    [Content]

    ## [Main Topic 2]
    [Content]

    ...

    ## Conclusion
    [Summary of key points]
    [Preview of next module if applicable]
    """
    response = get_gemini_response(input_text, prompt)
    return response

def generate_quiz(module_content):
    prompt = f"""
    Based on the following module content, generate a quiz with 3-5 multiple-choice questions (MCQs):

    {module_content}

    For each question:
    1. Provide the question
    2. List four possible answers (A, B, C, D)
    3. Indicate the correct answer

    Format the quiz as follows:
    # Quiz

    1. [Question 1]
       A. [Option A]
       B. [Option B]
       C. [Option C]
       D. [Option D]
    
    2. [Question 2]
       A. [Option A]
       B. [Option B]
       C. [Option C]
       D. [Option D]

    ...

    Answers:
    1. [Correct Answer]
    2. [Correct Answer]
    ...

    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    response = model.predict(prompt)
    return response

def convert_to_html(content):
    prompt = f"""
    Convert the following content into HTML format suitable for a lecture note:

    {content}

    Please structure the HTML document with proper headings, paragraphs, and formatting. 
    Include a title and author section. Ensure that any mathematical equations, lists, or 
    special formatting are properly converted to HTML syntax.
    Maintain the heading structure from the original content.
    Use CSS for basic styling to make the notes visually appealing.
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    response = model.predict(prompt)
    return response

def convert_to_latex(content):
    prompt = f"""
    Convert the following content into LaTeX format suitable for a lecture note:

    {content}

    Please structure the LaTeX document with proper sections, subsections, and formatting. 
    Include a title and author section. Ensure that any mathematical equations, lists, or 
    special formatting are properly converted to LaTeX syntax.
    Maintain the heading structure from the original content.
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    response = model.predict(prompt)
    return response

def convert_to_latex(content):
    prompt = f"""
    Convert the following content into LaTeX format suitable for a lecture note:

    {content}

    Please structure the LaTeX document with proper sections, subsections, and formatting. 
    Include a title and author section. Ensure that any mathematical equations, lists, or 
    special formatting are properly converted to LaTeX syntax.
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    response = model.predict(prompt)
    return response

def create_pdf_from_latex(latex_content):
    with tempfile.TemporaryDirectory() as tmpdir:
        tex_file = os.path.join(tmpdir, "notes.tex")
        with open(tex_file, "w") as f:
            f.write(latex_content)
        
        
        for _ in range(2):
            subprocess.run(["pdflatex", "-output-directory", tmpdir, tex_file], 
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        
        pdf_file = os.path.join(tmpdir, "notes.pdf")
        with open(pdf_file, "rb") as f:
            pdf_content = f.read()
    
    return pdf_content

def get_gemini_response(input_text, prompt):
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.7)
    prompt_template = PromptTemplate(
        input_variables=["input"],
        template=prompt
    )
    final_prompt = prompt_template.format(input=input_text)
    response = model.predict(final_prompt)
    return response

def generate_course_titles(profession, target_audience, course_content, course_description):
    input_text = f"""
    Profession: {profession}
    Target Audience: {target_audience}
    Course Content: {course_content}
    Course Description: {course_description}
    """
    prompt = """
    Based on the following information, generate 5 engaging and unique course titles:
    {input}

    Return the titles as a numbered list.
    """
    response = get_gemini_response(input_text, prompt)
    return [title.strip() for title in response.split("\n") if title.strip()]

def generate_course_outline(title, profession, target_audience, course_content, course_description):
    input_text = f"""
    Course Title: {title}
    Profession: {profession}
    Target Audience: {target_audience}
    Course Content: {course_content}
    Course Description: {course_description}
    """
    prompt = """
    Based on the following information, create a course outline with atmost 8 modules:
    {input}

    Generate a list of module names, each on a new line. Do not include any additional information or numbering.
    """
    response = get_gemini_response(input_text, prompt)
    return [module.strip() for module in response.split("\n") if module.strip()][:8]

def sarvam_api_call(endpoint, data):
    headers = {
        "Authorization": f"Bearer {SARVAM_API_KEY}",
        "Content-Type": "application/json"
    }
    response = requests.post(f"{SARVAM_API_URL}/{endpoint}", headers=headers, json=data)
    return response.json()

def perform_translation(text, source_lang="en-IN", target_lang="hi-IN"):
    url = f"{SARVAM_API_URL}/translate"
    
    payload = {
        "source_language_code": source_lang,
        "target_language_code": target_lang,
        "speaker_gender": "Female",
        "mode": "formal",
        "model": "mayura:v1",
        "enable_preprocessing": True,
        "input": text
    }
    headers = {
        "api-subscription-key": SARVAM_API_KEY,
        "Content-Type": "application/json"
    }

    response = requests.post(url, json=payload, headers=headers)
    
    if response.status_code == 200:
        return response.json()
    else:
        st.error(f"Translation failed with status code: {response.status_code}")
        st.error(f"Error message: {response.text}")
        return None

def perform_tts(text, speaker="meera"):
    url = f"{SARVAM_API_URL}/text-to-speech"
    
    payload = {
        "inputs": [text[:400]],
        "target_language_code": "hi-IN",
        "speaker": "meera",
        "pitch": 0,
        "pace": 1.65,
        "loudness": 1.5,
        "speech_sample_rate": 8000,
        "enable_preprocessing": True,
        "model": "bulbul:v1"
    }
    headers = {
        "api-subscription-key": SARVAM_API_KEY,
        "Content-Type": "application/json"
    }

    response = requests.post(url, json=payload, headers=headers)
    print(response.text)
    if response.status_code == 200:
        aud_str=response.text[12:-3]
        audio_data=base64.b64decode(aud_str)
        return audio_data
    else:
        st.error(f"TTS failed with status code: {response.status_code}")
        st.error(f"Error message: {response.text}")
        return None

def main():
    st.set_page_config(page_title="Course Creation Copilot", page_icon="🎓", layout="wide")
    st.title("SarvamScholar: Course Creation Copilot 🎓")

    if 'step' not in st.session_state:
        st.session_state.step = 'input'
    if 'course_titles' not in st.session_state:
        st.session_state.course_titles = []
    if 'selected_title' not in st.session_state:
        st.session_state.selected_title = ''
    if 'course_outline' not in st.session_state:
        st.session_state.course_outline = []
    if 'course_info' not in st.session_state:
        st.session_state.course_info = {
            'profession': '',
            'target_audience': '',
            'course_content': '',
            'course_description': ''
        }
    if 'module_content' not in st.session_state:
        st.session_state.module_contents = {}
    if 'current_module_index' not in st.session_state:
        st.session_state.current_module_index = 0
    if 'latex_contents' not in st.session_state:
        st.session_state.latex_contents = {}
    if 'quizzes' not in st.session_state:
        st.session_state.quizzes = {}
    if 'html_contents' not in st.session_state:
        st.session_state.html_contents = {}
    if 'ppt_contents' not in st.session_state:
        st.session_state.ppt_contents = {}

    st.sidebar.title("Navigation")
    step = st.sidebar.selectbox(
        "Go to step:",
        ['Input Information', 'Select Course Title', 'Customize Outline', 'Generate Content'],
        index=['input', 'select_title', 'customize_outline', 'generate_content'].index(st.session_state.step)
    )

    if st.sidebar.button("Back"):
        if st.session_state.step == 'select_title':
            st.session_state.step = 'input'
        elif st.session_state.step == 'customize_outline':
            st.session_state.step = 'select_title'
        elif st.session_state.step == 'generate_content':
            st.session_state.step = 'customize_outline'

    if step == 'Input Information':
        st.session_state.step = 'input'
    elif step == 'Select Course Title':
        st.session_state.step = 'select_title'
    elif step == 'Customize Outline':
        st.session_state.step = 'customize_outline'
    elif step == 'Generate Content':
        st.session_state.step = 'generate_content'

    if st.session_state.step == 'input':
        st.header("Step 1: Input Course Information")
        st.session_state.course_info['profession'] = st.text_input("What is your profession?", value=st.session_state.course_info['profession'])
        st.session_state.course_info['target_audience'] = st.text_input("Who is your target audience?", value=st.session_state.course_info['target_audience'])
        st.session_state.course_info['course_content'] = st.text_area("Describe the course content:", value=st.session_state.course_info['course_content'])
        st.session_state.course_info['course_description'] = st.text_area("Provide a brief course description: [opt]", value=st.session_state.course_info['course_description'])

        if st.button("Generate Course Titles"):
            with st.spinner("Generating course titles..."):
                st.session_state.course_titles = generate_course_titles(
                    st.session_state.course_info['profession'],
                    st.session_state.course_info['target_audience'],
                    st.session_state.course_info['course_content'],
                    st.session_state.course_info['course_description']
                )
                st.session_state.step = 'select_title'

    elif st.session_state.step == 'select_title':
        st.header("Step 2: Select Course Title")
        st.write("Choose a course title from the options below:")
        for i, title in enumerate(st.session_state.course_titles):
            if st.button(f"{i+1}. {title}"):
                st.session_state.selected_title = title
                with st.spinner("Generating course outline..."):
                    st.session_state.course_outline = generate_course_outline(
                        title,
                        st.session_state.course_info['profession'],
                        st.session_state.course_info['target_audience'],
                        st.session_state.course_info['course_content'],
                        st.session_state.course_info['course_description']
                    )
                st.session_state.step = 'customize_outline'


    elif st.session_state.step == 'customize_outline':
        st.header("Step 3: Customize Course Outline")
        st.subheader(st.session_state.selected_title)
        st.write("Drag and drop to reorder modules or add your own:")
        
        new_module = st.text_input("Add a new module:")
        if st.button("Add Module"):
            if new_module:
                st.session_state.course_outline.append(new_module)

        updated_outline = st.session_state.course_outline.copy()
        for i, module in enumerate(st.session_state.course_outline):
            col1, col2 = st.columns([0.9, 0.1])
            with col1:
                updated_module = st.text_input(f"Module {i+1}", value=module, key=f"module_{i}")
                updated_outline[i] = updated_module
            with col2:
                if st.button("🗑️", key=f"delete_{i}"):
                    updated_outline.pop(i)
        
        st.session_state.course_outline = updated_outline

        if st.button("Finalize Modules and Generate Content"):
            st.session_state.step = 'generate_content'
            st.session_state.current_module_index = 0

    elif st.session_state.step == 'generate_content':
        st.header("Step 4: Generate Content for Modules")
        st.subheader(st.session_state.selected_title)

        if st.session_state.current_module_index < len(st.session_state.course_outline):
            current_module = st.session_state.course_outline[st.session_state.current_module_index]
            st.write(f"Current Module: {current_module}")

            if current_module not in st.session_state.module_contents:
                with st.spinner(f"Generating content for module: {current_module}"):
                    previous_module = st.session_state.course_outline[st.session_state.current_module_index - 1] if st.session_state.current_module_index > 0 else None
                    module_content = generate_module_content(
                        st.session_state.selected_title,
                        current_module,
                        st.session_state.course_info,
                        previous_module
                    )
                    st.session_state.module_contents[current_module] = module_content
            st.write("Generated Content:")
            st.markdown(st.session_state.module_contents[current_module])
            if current_module not in st.session_state.quizzes:
                if st.button("Generate Quiz"):
                    with st.spinner("Generating quiz..."):
                        quiz_content = generate_quiz(st.session_state.module_contents[current_module])
                        st.session_state.quizzes[current_module] = quiz_content
                        st.success("Quiz generated successfully!")

            if current_module in st.session_state.quizzes:
                st.subheader("Module Quiz")
                st.markdown(st.session_state.quizzes[current_module])
                
            # if current_module not in st.session_state.latex_contents:
            #     if st.button("Generate LaTeX Notes"):
            #         with st.spinner("Generating LaTeX content..."):
            #             latex_content = convert_to_latex(st.session_state.module_contents[current_module])
            #             st.session_state.latex_contents[current_module] = latex_content
            #             st.success("LaTeX content generated successfully!")

            # if current_module in st.session_state.latex_contents:
            #     if st.button("Download Notes (PDF)"):
            #         pdf_content = create_pdf_from_latex(st.session_state.latex_contents[current_module])
            #         st.download_button(
            #             label="Download PDF",
            #             data=pdf_content,
            #             file_name=f"{current_module}_notes.pdf",
            #             mime="application/pdf"
            #         )
            if current_module not in st.session_state.html_contents:
                if st.button("Generate HTML Notes"):
                    with st.spinner("Generating HTML content..."):
                        full_content = f"{st.session_state.module_contents[current_module]}\n\n{st.session_state.quizzes.get(current_module, '')}"
                        html_content = convert_to_html(full_content)
                        st.session_state.html_contents[current_module] = html_content
                        st.success("HTML content generated successfully!")

            if current_module in st.session_state.html_contents:
                if st.download_button(
                    label="Download Notes (HTML)",
                    data=st.session_state.html_contents[current_module],
                    file_name=f"{current_module}_notes.html",
                    mime="text/html"
                ):
                    st.success("HTML notes downloaded successfully!")
            # if current_module not in st.session_state.latex_contents:
            #     if st.button("Generate LaTeX Notes"):
            #         with st.spinner("Generating LaTeX content..."):
            #             full_content = f"{st.session_state.module_contents[current_module]}\n\n{st.session_state.quizzes.get(current_module, '')}"
            #             latex_content = convert_to_latex(full_content)
            #             st.session_state.latex_contents[current_module] = latex_content
            #             st.success("LaTeX content generated successfully!")

            # if current_module in st.session_state.latex_contents:
            #     if st.button("Download Notes (PDF)"):
            #         pdf_content = create_pdf_from_latex(st.session_state.latex_contents[current_module])
            #         st.download_button(
            #             label="Download PDF",
            #             data=pdf_content,
            #             file_name=f"{current_module}_notes.pdf",
            #             mime="application/pdf"
            #         )
            # st.write("Generated Content:")
            # st.write(st.session_state.module_contents[current_module])
            st.subheader("PowerPoint Generation")
            if current_module not in st.session_state.ppt_contents:
                if st.button("Generate PowerPoint"):
                    with st.spinner("Generating PowerPoint..."):
                        ppt_content = create_ppt(st.session_state.module_contents[current_module], current_module)
                        st.session_state.ppt_contents[current_module] = ppt_content
                        st.success("PowerPoint generated successfully!")

            if current_module in st.session_state.ppt_contents:
                if st.download_button(
                    label="Download PowerPoint",
                    data=st.session_state.ppt_contents[current_module],
                    file_name=f"{current_module}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                ):
                    st.success("PowerPoint downloaded successfully!")
            st.subheader("Audio Processing Options")
            option = st.selectbox(
                "Choose an audio processing option:",
                ["None", "Translation to Hindi", "TTS (Text-to-Speech)"]
            )

            if option == "Translation to Hindi":
                if st.button("Translate to Hindi"):
                    with st.spinner("Translating..."):
                        result = perform_translation(st.session_state.module_contents[current_module])
                        if result:
                            st.write("Hindi Translation:")
                            st.write(result['translated_text'])
                            st.subheader("Audio Processing Options for Hindi")
                            optionh = st.selectbox(
                                "Text to Speech for Hindi Content:",
                                ["None", "TTS (Hindi-to-Speech)","Generate Hindi PowerPoint"]
                            )
                            if optionh == "TTS (Hindi-to-Speech)":
                                speakerh = st.selectbox("Choose a speaker:", ["meera", "pavithra", "maitreyi", "arvind", "amol", "amartya"])
                                if st.button("Generate Hindi Speech"):
                                    with st.spinner("Generating speech..."):
                                        audio_datah = perform_tts(result['translated_text'], speakerh)
                                        if audio_datah:
                                            with open("t2sh.wav", "wb") as wav_audioh:
                                                wav_audioh.write(audio_datah)
                                            with open("t2sh.wav", "rb") as wav_audio_file:
                                                audio_bytesh = wav_audio_file.read()
                                            st.audio(audio_bytesh, format="audio/wav")
                                        else:
                                            st.error("Speech generation failed. Please try again.")
                            elif optionh == "Generate Hindi PowerPoint":
                                    if st.button("Generate Hindi PowerPoint"):
                                        with st.spinner("Generating Hindi PowerPoint..."):
                                            hindi_ppt_content = create_ppt(result['translated_text'], f"{current_module} (Hindi)", is_hindi=True)
                                            st.session_state.ppt_contents[f"{current_module}_hindi"] = hindi_ppt_content
                                            st.success("Hindi PowerPoint generated successfully!")
                                    
                                    if f"{current_module}_hindi" in st.session_state.ppt_contents:
                                        if st.download_button(
                                            label="Download Hindi PowerPoint",
                                            data=st.session_state.ppt_contents[f"{current_module}_hindi"],
                                            file_name=f"{current_module}_hindi.pptx",
                                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                        ):
                                            st.success("Hindi PowerPoint downloaded successfully!")
                        else:
                            st.error("Translation failed. Please try again.")
            elif option == "TTS (Text-to-Speech)":
                speaker = st.selectbox("Choose a speaker:", ["meera", "pavithra", "maitreyi", "arvind", "amol", "amartya"])
                if st.button("Generate Speech"):
                    with st.spinner("Generating speech..."):
                        audio_data = perform_tts(st.session_state.module_contents[current_module], speaker)
                        if audio_data:
                            with open("t2s.wav", "wb") as wav_audio:
                                wav_audio.write(audio_data)
                            with open("t2s.wav", "rb") as wav_audio_file:
                                audio_bytes = wav_audio_file.read()
                            st.audio(audio_bytes, format="audio/wav")
                        else:
                            st.error("Speech generation failed. Please try again.")

            if st.button("Next Module"):
                st.session_state.current_module_index += 1
        else:
            st.write("All modules have been processed.")

if __name__ == "__main__":
    main()
